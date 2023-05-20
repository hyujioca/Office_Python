from pptx import Presentation
prs = Presentation("test.pptx")
title_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hey guys!"

text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
subtitle.text = "文字数: {}".format(text_runs)
prs.save("test2.pptx")